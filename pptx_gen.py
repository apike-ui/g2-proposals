"""
G2 Proposal Builder — PPTX generator.
Professional, G2-branded 16:9 PowerPoint deck with comprehensive pricing detail.
"""

import io, os, math
from datetime import date, datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ────────────────────────────────────────────────────
RORANGE  = RGBColor(0xFF, 0x49, 0x2C)
NAVY     = RGBColor(0x06, 0x28, 0x46)
GREEN    = RGBColor(0x27, 0xD3, 0xBC)
BLUE     = RGBColor(0x00, 0x73, 0xF5)
PURPLE   = RGBColor(0x57, 0x46, 0xB2)
YELLOW   = RGBColor(0xFF, 0xC8, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xF2, 0xF4, 0xF7)
MID      = RGBColor(0x4D, 0x64, 0x80)
GREY     = RGBColor(0xD9, 0xDE, 0xE6)
NAVY_D   = RGBColor(0x04, 0x1D, 0x35)
GREEN_L  = RGBColor(0xE0, 0xFA, 0xF6)
BLUE_L   = RGBColor(0xE0, 0xEE, 0xFF)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Figtree"
_HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(_HERE, "assets")


# ── Helpers ──────────────────────────────────────────────────────────

def _r(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def _t(slide, x, y, w, h, text, sz=12, bold=False, color=NAVY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    return tb

def _logo(slide, x, y, h, variant="inverse"):
    names = {"inverse": "G2Logo-Inverse.png", "red": "G2Logo-Red.png"}
    path = os.path.join(ASSETS, names.get(variant, "G2Logo-Inverse.png"))
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))

def _hdr(slide, title, subtitle=""):
    _r(slide, 0, 0, 13.333, 1.2, NAVY)
    _r(slide, 0, 0, 0.08, 1.2, RORANGE)
    _logo(slide, 11.8, 0.25, 0.55)
    _t(slide, 0.7, 0.12, 8, 0.5, title, sz=26, bold=True, color=WHITE)
    if subtitle:
        _t(slide, 0.7, 0.68, 8, 0.3, subtitle, sz=12, color=GREY)
    _r(slide, 0, 1.2, 13.333, 0.05, RORANGE)

def _ftr(slide, page=None):
    _r(slide, 0, 7.1, 13.333, 0.4, NAVY)
    _t(slide, 0.5, 7.14, 3, 0.28, "g2.com", sz=9, bold=True, color=WHITE)
    _logo(slide, 12.2, 7.14, 0.24)
    if page:
        _t(slide, 10, 7.14, 2, 0.28, str(page), sz=9, color=GREY, align=PP_ALIGN.RIGHT)

def _money(v):
    try:
        n = float(v)
        return f"${int(n):,}" if n == int(n) else f"${n:,.0f}"
    except:
        return str(v)

def _fdate(ds):
    try:
        return datetime.strptime(ds, "%Y-%m-%d").strftime("%b %d, %Y")
    except:
        return ds or ""


# ── Slide 1: Cover ──────────────────────────────────────────────────

def _cover(prs, data, ct):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _r(sl, 0, 0, 13.333, 7.5, NAVY)
    _r(sl, 0, 0, 0.08, 7.5, RORANGE)
    _logo(sl, 11.3, 0.4, 0.7)

    _t(sl, 0.7, 0.5, 4, 0.3, "PROPOSAL", sz=13, bold=True, color=RORANGE)
    cust = data.get("cust") or "Your Company"
    _t(sl, 0.7, 1.3, 11, 1.5, cust, sz=48, bold=True, color=WHITE, wrap=True)
    _r(sl, 0.7, 3.2, 4, 0.04, RORANGE)
    _t(sl, 0.7, 3.5, 9, 0.4, "Custom Investment Summary  |  G2 Buyer Intelligence Platform", sz=15, color=GREY)

    iy = 4.2
    rep = data.get("rep") or ""
    today = data.get("date") or date.today().strftime("%B %d, %Y")
    term = data.get("contractTerm") or "12"
    sd = _fdate(data.get("startDate"))
    ed = _fdate(data.get("endDate"))

    if rep:
        _t(sl, 0.7, iy, 6, 0.3, f"Prepared by  {rep}", sz=13, color=MID); iy += 0.35
    _t(sl, 0.7, iy, 6, 0.3, today, sz=13, color=MID); iy += 0.35
    term_s = f"{term}-month term" if term and term != "custom" else "Custom term"
    _t(sl, 0.7, iy, 6, 0.3, term_s, sz=13, color=MID); iy += 0.35
    if sd or ed:
        _t(sl, 0.7, iy, 6, 0.3, f"{sd or 'TBD'}  —  {ed or 'TBD'}", sz=13, bold=True, color=BLUE)

    # Stats cards
    cx, cw = 9.0, 3.8
    for i, (stat, desc, ac) in enumerate([
        ("90M+", "annual buyer interactions", RORANGE),
        ("#1", "software marketplace", GREEN),
        ("2.6M+", "verified reviews", YELLOW),
        ("160K+", "products & services", BLUE),
    ]):
        cy = 1.3 + i * 1.1
        _r(sl, cx, cy, cw, 0.9, NAVY_D)
        _r(sl, cx, cy, 0.06, 0.9, ac)
        _t(sl, cx + 0.2, cy + 0.1, 1.5, 0.4, stat, sz=20, bold=True, color=ac)
        _t(sl, cx + 1.8, cy + 0.15, 1.9, 0.55, desc, sz=10, color=GREY, wrap=True)

    _r(sl, 0, 6.95, 13.333, 0.55, RORANGE)
    _t(sl, 0.7, 7.0, 8, 0.35, f"Prepared exclusively for {cust}", sz=10, bold=True, color=WHITE)
    _logo(sl, 12.0, 7.0, 0.35)


# ── Slide 2: Executive Summary ──────────────────────────────────────

def _exec_summary(prs, data, ct, pg):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _r(sl, 0, 0, 13.333, 7.5, WHITE)
    cust = data.get("cust") or ""
    _hdr(sl, "Investment Summary", cust)

    grand = ct.get("grandTotal") or 0
    acv = ct.get("finalAcv") or 0
    non_acv = ct.get("totalNonAcv") or 0
    n_prod = len(ct.get("prodSummaries") or [])
    savings = ct.get("totalList", 0) - acv
    term = int(data.get("contractTerm") or 12)

    # Stat cards
    cards = [
        ("Total ACV", _money(acv), BLUE),
        ("Non-ACV", _money(non_acv), PURPLE),
        ("Products", str(n_prod), NAVY),
        ("Total Investment", _money(grand) + "/yr", GREEN),
    ]
    cw, gap = 2.65, 0.3
    sx = (13.333 - (4 * cw + 3 * gap)) / 2
    for i, (lbl, val, bg) in enumerate(cards):
        cx = sx + i * (cw + gap)
        _r(sl, cx, 1.55, cw, 1.7, bg)
        _t(sl, cx, 1.7, cw, 0.3, lbl, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _t(sl, cx, 2.1, cw, 0.8, val, sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # TCV callout for multi-year
    if term > 12:
        years = math.ceil(term / 12)
        tcv = grand * years
        _r(sl, sx, 3.4, 4 * cw + 3 * gap, 0.4, BLUE_L)
        _r(sl, sx, 3.4, 0.05, 0.4, BLUE)
        _t(sl, sx + 0.2, 3.46, 6, 0.28, f"Total Contract Value ({years} years): {_money(tcv)}", sz=12, bold=True, color=NAVY)

    # Product breakdown table
    ty = 4.1 if term > 12 else 3.7
    _t(sl, 0.7, ty, 10, 0.3, "SKU Summary by Profile", sz=13, bold=True, color=NAVY)
    _r(sl, 0.7, ty + 0.35, 2, 0.04, RORANGE)

    tw = 11.9
    mx = 0.7
    rh = 0.32
    # Header
    hy = ty + 0.55
    _r(sl, mx, hy, tw, 0.35, NAVY)
    cols = [("Profile", 3.5, PP_ALIGN.LEFT), ("Package", 1.5, PP_ALIGN.CENTER),
            ("Add-Ons", 4.0, PP_ALIGN.LEFT), ("ACV", 1.5, PP_ALIGN.RIGHT)]
    cx = mx
    for lbl, w, al in cols:
        _t(sl, cx + 0.1, hy + 0.04, w - 0.2, 0.25, lbl, sz=10, bold=True, color=WHITE, align=al)
        cx += w
    # Product rows
    ry = hy + 0.35
    prods = ct.get("prodSummaries") or []
    for j, ps in enumerate(prods[:7]):
        bg = LIGHT if j % 2 == 0 else WHITE
        _r(sl, mx, ry, tw, rh, bg)
        cx = mx
        _t(sl, cx + 0.1, ry + 0.04, 3.3, rh - 0.08, ps.get("name", ""), sz=10, bold=True, color=NAVY)
        cx += 3.5
        pkg = ps.get("basePkg")
        pkg_name = pkg.get("name", "") if isinstance(pkg, dict) else str(pkg or "").title()
        _t(sl, cx + 0.1, ry + 0.04, 1.3, rh - 0.08, pkg_name, sz=10, color=MID, align=PP_ALIGN.CENTER)
        cx += 1.5
        addons = ps.get("addonLines") or []
        addon_str = ", ".join([a.get("name", "") + (f" x{a['qty']}" if a.get("qty", 1) > 1 else "") for a in addons]) or "—"
        _t(sl, cx + 0.1, ry + 0.04, 3.8, rh - 0.08, addon_str, sz=9, color=MID)
        cx += 4.0
        _t(sl, cx + 0.1, ry + 0.04, 1.3, rh - 0.08, _money(ps.get("prodAcv", 0)), sz=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        ry += rh

    if len(prods) > 7:
        _t(sl, mx, ry, tw, 0.25, f"+ {len(prods) - 7} more (see detail slides)", sz=9, color=MID, align=PP_ALIGN.CENTER)

    # Non-ACV row
    non_acv_lines = ct.get("nonAcvLines") or []
    if non_acv_lines:
        ry += 0.05
        _r(sl, mx, ry, tw, rh, LIGHT)
        names = ", ".join([l.get("name", "") for l in non_acv_lines])
        _t(sl, mx + 0.1, ry + 0.04, 8, rh - 0.08, f"Non-ACV: {names}", sz=9, color=MID)
        _t(sl, mx + tw - 1.6, ry + 0.04, 1.4, rh - 0.08, _money(non_acv), sz=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    _ftr(sl, pg)


# ── Slide 3+: Per-product detail ────────────────────────────────────

def _product_detail(prs, data, ps, pg):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _r(sl, 0, 0, 13.333, 7.5, WHITE)
    name = ps.get("name") or "Product"
    pkg = ps.get("basePkg")
    pkg_name = pkg.get("name", "") if isinstance(pkg, dict) else str(pkg or "").title()
    _hdr(sl, name, f"{pkg_name} Package")

    # Left: pricing breakdown
    lx, lw = 0.7, 6.0

    # Base package
    _t(sl, lx, 1.55, lw, 0.3, "BASE PACKAGE", sz=10, bold=True, color=MID)
    _r(sl, lx, 1.85, lw, 0.55, LIGHT)
    _r(sl, lx, 1.85, 0.05, 0.55, PURPLE)
    _t(sl, lx + 0.2, 1.92, 3.5, 0.35, f"{pkg_name}", sz=13, bold=True, color=NAVY)
    base_rate = ps.get("baseRate", 0)
    list_base = ps.get("listBase", 0)
    _t(sl, lx + lw - 1.8, 1.92, 1.6, 0.35, f"{_money(base_rate)}/yr", sz=13, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    if list_base and list_base != base_rate and list_base > 0:
        disc = round((1 - base_rate / list_base) * 100)
        _t(sl, lx + lw - 1.8, 2.22, 1.6, 0.2, f"List: {_money(list_base)} ({disc}% off)", sz=9, color=GREEN, align=PP_ALIGN.RIGHT)

    # Add-ons
    addons = ps.get("addonLines") or []
    if addons:
        _t(sl, lx, 2.6, lw, 0.3, "ADD-ON PRODUCTS", sz=10, bold=True, color=MID)
        _r(sl, lx, 2.88, lw, 0.04, GREY)

        # Table header
        _r(sl, lx, 2.95, lw, 0.3, NAVY)
        _t(sl, lx + 0.1, 2.98, 2.2, 0.22, "SKU", sz=9, bold=True, color=WHITE)
        _t(sl, lx + 2.3, 2.98, 1.5, 0.22, "Tier / Detail", sz=9, bold=True, color=WHITE)
        _t(sl, lx + 3.8, 2.98, 0.6, 0.22, "Qty", sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _t(sl, lx + 4.4, 2.98, 0.7, 0.22, "Rate", sz=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        _t(sl, lx + 5.1, 2.98, 0.8, 0.22, "Total", sz=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

        ay = 3.25
        for k, a in enumerate(addons[:10]):
            bg = LIGHT if k % 2 == 0 else WHITE
            _r(sl, lx, ay, lw, 0.35, bg)
            _t(sl, lx + 0.1, ay + 0.05, 2.1, 0.25, a.get("name", ""), sz=10, bold=True, color=NAVY)
            _t(sl, lx + 2.3, ay + 0.05, 1.5, 0.25, a.get("tier", ""), sz=9, color=MID)
            qty = a.get("qty", 1)
            _t(sl, lx + 3.8, ay + 0.05, 0.6, 0.25, str(qty), sz=9, color=MID, align=PP_ALIGN.CENTER)
            _t(sl, lx + 4.4, ay + 0.05, 0.7, 0.25, _money(a.get("rate", 0)), sz=9, color=MID, align=PP_ALIGN.RIGHT)
            _t(sl, lx + 5.1, ay + 0.05, 0.8, 0.25, _money(a.get("total", 0)), sz=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
            # Discount indicator
            disc = a.get("disc", 0)
            if disc and disc > 0:
                _t(sl, lx + 0.1, ay + 0.22, 2, 0.12, f"{disc}% off list", sz=7, color=GREEN)
            ay += 0.35

    # Right: Product ACV card
    rx, rw = 7.2, 5.5
    _r(sl, rx, 1.55, rw, 2.0, NAVY)
    _r(sl, rx, 1.55, rw, 0.06, RORANGE)
    _t(sl, rx + 0.3, 1.8, rw - 0.6, 0.3, "PRODUCT ACV", sz=10, bold=True, color=GREY)
    _t(sl, rx + 0.3, 2.2, rw - 0.6, 0.7, _money(ps.get("prodAcv", 0)), sz=40, bold=True, color=WHITE)
    _t(sl, rx + 0.3, 2.9, rw - 0.6, 0.3, "/year", sz=14, color=GREY)

    # Breakdown
    _t(sl, rx, 3.8, rw, 0.25, "Investment Breakdown", sz=12, bold=True, color=NAVY)
    _r(sl, rx, 4.05, rw, 0.03, RORANGE)
    by = 4.2
    _r(sl, rx, by, rw, 0.35, LIGHT)
    _t(sl, rx + 0.15, by + 0.05, 3.5, 0.25, f"Base: {pkg_name}", sz=10, color=NAVY)
    _t(sl, rx + rw - 1.5, by + 0.05, 1.3, 0.25, _money(base_rate), sz=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    by += 0.35
    for a in addons:
        bg = WHITE if (by - 4.2) / 0.35 % 2 < 1 else LIGHT
        _r(sl, rx, by, rw, 0.35, bg)
        _t(sl, rx + 0.15, by + 0.05, 3.5, 0.25, a.get("name", ""), sz=10, color=NAVY)
        _t(sl, rx + rw - 1.5, by + 0.05, 1.3, 0.25, _money(a.get("total", 0)), sz=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        by += 0.35
        if by > 6.5:
            break
    # Total
    by += 0.1
    _r(sl, rx, by, rw, 0.4, NAVY)
    _t(sl, rx + 0.15, by + 0.06, 3, 0.28, "Total Product ACV", sz=11, bold=True, color=WHITE)
    _t(sl, rx + rw - 1.8, by + 0.06, 1.6, 0.28, _money(ps.get("prodAcv", 0)), sz=14, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

    _ftr(sl, pg)


# ── Pricing table slide ─────────────────────────────────────────────

def _pricing_table(prs, data, ct, pg):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _r(sl, 0, 0, 13.333, 7.5, WHITE)
    _hdr(sl, "Pricing Details", data.get("cust") or "")

    mx, tw = 0.7, 11.9
    rh = 0.34

    # Header
    hy = 1.5
    _r(sl, mx, hy, tw, 0.38, NAVY)
    _t(sl, mx + 0.15, hy + 0.06, 5.5, 0.26, "Item", sz=10, bold=True, color=WHITE)
    _t(sl, mx + 5.7, hy + 0.06, 1.5, 0.26, "Qty", sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _t(sl, mx + 7.2, hy + 0.06, 1.5, 0.26, "Rate", sz=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    _t(sl, mx + 8.7, hy + 0.06, 1.5, 0.26, "List", sz=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    _t(sl, mx + 10.2, hy + 0.06, 1.5, 0.26, "Total", sz=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    ry = hy + 0.38
    prods = ct.get("prodSummaries") or []

    for ps in prods:
        if ry > 5.8:
            break
        name = ps.get("name", "Product")
        pkg = ps.get("basePkg")
        pkg_name = pkg.get("name", "") if isinstance(pkg, dict) else str(pkg or "").title()

        # Product header row
        _r(sl, mx, ry, tw, rh, BLUE_L)
        _r(sl, mx, ry, 0.05, rh, BLUE)
        _t(sl, mx + 0.2, ry + 0.04, 6, rh - 0.08, f"{name} — {pkg_name}", sz=10, bold=True, color=NAVY)
        _t(sl, mx + 10.2, ry + 0.04, 1.5, rh - 0.08, _money(ps.get("prodAcv", 0)), sz=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        ry += rh

        # Base row
        base = ps.get("baseRate", 0)
        list_b = ps.get("listBase", 0)
        if base > 0:
            bg = LIGHT
            _r(sl, mx, ry, tw, rh, bg)
            _t(sl, mx + 0.4, ry + 0.04, 5.3, rh - 0.08, f"Base: {pkg_name}", sz=9, color=MID)
            _t(sl, mx + 5.7, ry + 0.04, 1.5, rh - 0.08, "1", sz=9, color=MID, align=PP_ALIGN.CENTER)
            _t(sl, mx + 7.2, ry + 0.04, 1.5, rh - 0.08, _money(base), sz=9, color=MID, align=PP_ALIGN.RIGHT)
            _t(sl, mx + 8.7, ry + 0.04, 1.5, rh - 0.08, _money(list_b) if list_b else "", sz=9, color=GREY, align=PP_ALIGN.RIGHT)
            _t(sl, mx + 10.2, ry + 0.04, 1.5, rh - 0.08, _money(base), sz=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
            ry += rh

        # Addon rows
        for a in (ps.get("addonLines") or []):
            if ry > 5.8:
                break
            bg = WHITE if int((ry - hy) / rh) % 2 == 0 else LIGHT
            _r(sl, mx, ry, tw, rh, bg)
            _t(sl, mx + 0.4, ry + 0.04, 3, rh - 0.08, a.get("name", ""), sz=9, color=NAVY)
            _t(sl, mx + 3.5, ry + 0.04, 2.2, rh - 0.08, a.get("tier", ""), sz=8, color=MID)
            qty = a.get("qty", 1)
            _t(sl, mx + 5.7, ry + 0.04, 1.5, rh - 0.08, str(qty), sz=9, color=MID, align=PP_ALIGN.CENTER)
            _t(sl, mx + 7.2, ry + 0.04, 1.5, rh - 0.08, _money(a.get("rate", 0)), sz=9, color=MID, align=PP_ALIGN.RIGHT)
            lt = a.get("listTotal") or a.get("listRate", 0)
            _t(sl, mx + 8.7, ry + 0.04, 1.5, rh - 0.08, _money(lt) if lt else "", sz=9, color=GREY, align=PP_ALIGN.RIGHT)
            _t(sl, mx + 10.2, ry + 0.04, 1.5, rh - 0.08, _money(a.get("total", 0)), sz=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
            ry += rh

    # Non-ACV
    non_acv = ct.get("nonAcvLines") or []
    if non_acv and ry < 5.8:
        ry += 0.08
        _r(sl, mx, ry, tw, rh, LIGHT)
        _r(sl, mx, ry, 0.05, rh, MID)
        _t(sl, mx + 0.2, ry + 0.04, 5, rh - 0.08, "Non-ACV Items", sz=10, bold=True, color=MID)
        _t(sl, mx + 10.2, ry + 0.04, 1.5, rh - 0.08, _money(ct.get("totalNonAcv", 0)), sz=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        ry += rh
        for na in non_acv:
            if ry > 5.8:
                break
            _r(sl, mx, ry, tw, rh, WHITE)
            _t(sl, mx + 0.4, ry + 0.04, 5, rh - 0.08, na.get("name", ""), sz=9, color=MID)
            qty = na.get("qty", 1)
            _t(sl, mx + 5.7, ry + 0.04, 1.5, rh - 0.08, str(qty), sz=9, color=MID, align=PP_ALIGN.CENTER)
            _t(sl, mx + 7.2, ry + 0.04, 1.5, rh - 0.08, _money(na.get("rate", 0)), sz=9, color=MID, align=PP_ALIGN.RIGHT)
            _t(sl, mx + 10.2, ry + 0.04, 1.5, rh - 0.08, _money(na.get("total", 0)), sz=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
            ry += rh

    # Totals
    ty = max(ry + 0.15, 5.6)
    _r(sl, mx, ty, tw, 0.03, NAVY)

    # Discounts
    disc_pct = ct.get("propDiscPct", 0)
    disc_amt = ct.get("propDiscAmt", 0)
    rc_disc = ct.get("rcDiscAmt", 0)
    rc_lines = ct.get("rcDiscLines") or []

    if disc_amt > 0:
        _r(sl, mx, ty + 0.08, tw, 0.3, GREEN_L)
        _r(sl, mx, ty + 0.08, 0.05, 0.3, GREEN)
        _t(sl, mx + 0.2, ty + 0.12, 6, 0.22, f"Proposal Discount ({disc_pct}%)", sz=10, bold=True, color=GREEN)
        _t(sl, mx + 10.2, ty + 0.12, 1.5, 0.22, f"-{_money(disc_amt)}", sz=10, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        ty += 0.35

    for dl in rc_lines:
        _r(sl, mx, ty + 0.08, tw, 0.3, BLUE_L)
        _r(sl, mx, ty + 0.08, 0.05, 0.3, BLUE)
        nm = dl.get("name", "Rate card discount")
        pct = dl.get("pct", 0)
        _t(sl, mx + 0.2, ty + 0.12, 6, 0.22, f"{nm} ({pct}%)", sz=10, bold=True, color=BLUE)
        _t(sl, mx + 10.2, ty + 0.12, 1.5, 0.22, f"-{_money(dl.get('amount', 0))}", sz=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        ty += 0.35

    # Grand total
    _r(sl, mx, ty + 0.08, tw, 0.5, NAVY)
    _t(sl, mx + 0.2, ty + 0.15, 6, 0.35, "Total Annual Investment", sz=14, bold=True, color=WHITE)
    _t(sl, mx + 10.2 - 1, ty + 0.15, 2.5, 0.35, _money(ct.get("grandTotal", 0)), sz=18, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

    savings = ct.get("totalList", 0) - ct.get("finalAcv", 0)
    if savings > 0:
        _t(sl, mx + 0.2, ty + 0.42, 6, 0.15, f"Total savings: {_money(savings)}", sz=8, color=GREEN)

    _ftr(sl, pg)


# ── Cost by Year ────────────────────────────────────────────────────

def _cost_by_year(prs, data, ct, pg):
    term = int(data.get("contractTerm") or 12)
    if term <= 12:
        return False
    years = math.ceil(term / 12)
    grand = ct.get("grandTotal") or 0

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _r(sl, 0, 0, 13.333, 7.5, WHITE)
    _hdr(sl, "Cost by Year", data.get("cust") or "")

    mx, tw = 0.7, 11.9
    hy = 1.6
    _r(sl, mx, hy, tw, 0.4, NAVY)
    for i, (lbl, al) in enumerate([("Year", PP_ALIGN.LEFT), ("Annual Total", PP_ALIGN.RIGHT), ("Cumulative", PP_ALIGN.RIGHT)]):
        w = 4.0 if i == 0 else 3.95
        _t(sl, mx + i * 4.0 + 0.15, hy + 0.06, w - 0.3, 0.26, lbl, sz=11, bold=True, color=WHITE, align=al)

    ry = hy + 0.4
    cum = 0
    for y in range(1, years + 1):
        partial = (y == years and term % 12 != 0)
        frac = (term % 12) / 12.0 if partial else 1.0
        yr_total = round(grand * frac)
        cum += yr_total
        bg = LIGHT if y % 2 == 1 else WHITE
        _r(sl, mx, ry, tw, 0.45, bg)
        lbl = f"Year {y}" + (f" ({term % 12}mo)" if partial else "")
        _t(sl, mx + 0.15, ry + 0.08, 3.7, 0.28, lbl, sz=12, bold=True, color=NAVY)
        _t(sl, mx + 4.15, ry + 0.08, 3.7, 0.28, _money(yr_total), sz=12, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        _t(sl, mx + 8.15, ry + 0.08, 3.6, 0.28, _money(cum), sz=12, color=NAVY, align=PP_ALIGN.RIGHT)
        ry += 0.45

    _r(sl, mx, ry + 0.1, tw, 0.5, NAVY)
    _t(sl, mx + 0.15, ry + 0.17, 5, 0.35, f"Total Contract Value ({years} year{'s' if years > 1 else ''})", sz=14, bold=True, color=WHITE)
    _t(sl, mx + 4.15, ry + 0.17, 3.7, 0.35, _money(cum), sz=16, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

    _ftr(sl, pg)
    return True


# ── Next Steps ──────────────────────────────────────────────────────

def _next_steps(prs, data, pg):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _r(sl, 0, 0, 13.333, 7.5, NAVY)
    _r(sl, 0, 0, 0.08, 7.5, RORANGE)
    _logo(sl, 11.3, 0.4, 0.55)

    _t(sl, 0.7, 0.5, 4, 0.3, "NEXT STEPS", sz=12, bold=True, color=RORANGE)
    _t(sl, 0.7, 1.0, 7.5, 0.7, "Let's Move Forward", sz=38, bold=True, color=WHITE)
    _r(sl, 0.7, 1.85, 3, 0.04, RORANGE)

    steps = [
        ("01", "Review This Proposal", "Share with stakeholders and confirm products and pricing align with your goals."),
        ("02", "Sign the Agreement", "We'll send the MSA and Order Form for e-signature via DocuSign."),
        ("03", "Onboarding Kickoff", "Your dedicated CSM will schedule onboarding within 5 business days."),
        ("04", "Go Live", "Access your G2 dashboard, integrate your CRM, and start capturing buyer intent."),
    ]
    colors = [RORANGE, GREEN, BLUE, PURPLE]
    for i, (num, title, desc) in enumerate(steps):
        sy = 2.2 + i * 1.05
        _r(sl, 0.7, sy, 0.6, 0.5, colors[i])
        _t(sl, 0.7, sy + 0.06, 0.6, 0.38, num, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _t(sl, 1.5, sy + 0.02, 6.5, 0.28, title, sz=14, bold=True, color=WHITE)
        _t(sl, 1.5, sy + 0.32, 6.5, 0.45, desc, sz=10, color=GREY, wrap=True)

    # Rep card
    cx, cw, cy = 8.8, 3.8, 1.0
    _r(sl, cx, cy, cw, 5.3, NAVY_D)
    _r(sl, cx, cy, cw, 0.05, RORANGE)
    _t(sl, cx + 0.3, cy + 0.3, cw - 0.6, 0.25, "YOUR G2 CONTACT", sz=10, bold=True, color=RORANGE)
    rep = data.get("rep") or "Your G2 Account Executive"
    _t(sl, cx + 0.3, cy + 0.7, cw - 0.6, 0.6, rep, sz=18, bold=True, color=WHITE, wrap=True)
    rep_title = data.get("repTitle") or "Account Executive"
    _t(sl, cx + 0.3, cy + 1.35, cw - 0.6, 0.25, rep_title, sz=11, color=GREY)
    _r(sl, cx + 0.3, cy + 1.7, cw - 0.6, 0.03, RORANGE)
    iy = cy + 1.9
    for field in ["repEmail", "repPhone"]:
        val = data.get(field)
        if val:
            _t(sl, cx + 0.3, iy, cw - 0.6, 0.25, val, sz=11, color=WHITE)
            iy += 0.3
    _logo(sl, cx + 0.3, cy + 3.5, 0.4)
    _t(sl, cx + 0.85, cy + 3.55, 2.5, 0.35, "The World's Largest\nSoftware Marketplace", sz=9, color=GREY, wrap=True)
    _t(sl, cx + 0.3, cy + 4.2, cw - 0.6, 0.8, "Questions? Reach out anytime.\nWe're committed to your success with G2.", sz=10, color=MID, wrap=True)

    _r(sl, 0, 6.95, 13.333, 0.55, RORANGE)
    cust = data.get("cust") or ""
    _t(sl, 0.7, 7.0, 8, 0.35, f"Prepared for {cust}  |  Confidential", sz=10, bold=True, color=WHITE)
    _logo(sl, 12.0, 7.0, 0.35)


# ── Main ─────────────────────────────────────────────────────────────

def build_pptx(data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Use pre-computed totals from frontend if available, else compute
    ct = data.get("_computed") or {}
    if not ct:
        ct = _fallback_compute(data)

    _cover(prs, data, ct)

    pg = 2
    _exec_summary(prs, data, ct, pg); pg += 1

    for ps in (ct.get("prodSummaries") or []):
        _product_detail(prs, data, ps, pg); pg += 1

    _pricing_table(prs, data, ct, pg); pg += 1

    if _cost_by_year(prs, data, ct, pg):
        pg += 1

    _next_steps(prs, data, pg)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _fallback_compute(data):
    """Basic fallback when frontend doesn't send _computed."""
    products = data.get("products") or []
    acct = data.get("acctItems") or {}
    disc = float(data.get("proposalDisc") or 0)
    summaries = []
    total_acv = 0
    total_list = 0
    for prod in products:
        pkg = (prod.get("basePkg") or "free").lower()
        list_b = {"free": 0, "professional": 18000, "enterprise": 36000}.get(pkg, 0)
        br = prod.get("baseRate")
        base = float(br) if br and float(br) > 0 else list_b
        p_acv = base
        p_list = list_b
        lines = []
        addons = prod.get("addons") or {}
        if isinstance(addons, dict):
            for aid, ap in addons.items():
                if isinstance(ap, dict):
                    if not ap.get("on", True): continue
                    q = int(ap.get("qty") or 1)
                    r = float(ap.get("rate") or 0)
                    t = r * q
                    lines.append({"name": ap.get("customDesc") or aid, "tier": "", "qty": q, "rate": r, "total": t, "listRate": 0, "listTotal": 0, "disc": 0})
                else:
                    t = float(ap) if ap else 0
                    lines.append({"name": aid, "tier": "", "qty": 1, "rate": t, "total": t, "listRate": t, "listTotal": t, "disc": 0})
                p_acv += t; p_list += t
        summaries.append({"name": prod.get("name") or "Product", "basePkg": pkg, "baseRate": base, "listBase": list_b, "addonLines": lines, "prodAcv": p_acv, "prodList": p_list})
        total_acv += p_acv; total_list += p_list
    non_acv = []
    t_na = 0
    for aid, ap in acct.items():
        if isinstance(ap, dict):
            if not ap.get("on", True): continue
            q = int(ap.get("qty") or 1); r = float(ap.get("rate") or 0); t = r * q
        else:
            t = float(ap) if ap else 0; q = 1; r = t
        non_acv.append({"name": aid, "qty": q, "rate": r, "total": t})
        t_na += t
    d_amt = round(total_acv * disc / 100) if disc > 0 else 0
    final = total_acv - d_amt
    return {"prodSummaries": summaries, "totalAcv": total_acv, "totalList": total_list,
            "nonAcvLines": non_acv, "totalNonAcv": t_na, "propDiscPct": disc, "propDiscAmt": d_amt,
            "rcDiscAmt": 0, "rcDiscLines": [], "finalAcv": final, "grandTotal": final + t_na}


if __name__ == "__main__":
    import sys
    sample = {
        "cust": "Acme Corp", "rep": "Jane Smith", "repEmail": "jane@g2.com",
        "repTitle": "Senior AE", "contractTerm": "24", "startDate": "2026-02-01", "endDate": "2028-01-31",
        "products": [
            {"name": "Sales Cloud", "basePkg": "enterprise", "baseRate": "30000",
             "addons": {"intent": {"on": True, "tierIdx": 3, "qty": 1, "rate": "25000"},
                        "content": {"on": True, "tierIdx": 0, "qty": 1, "rate": "40000"}}},
            {"name": "Service Cloud", "basePkg": "professional", "baseRate": "15000",
             "addons": {"rms": {"on": True, "tierIdx": 1, "qty": 1, "rate": "18000"}}},
        ],
        "acctItems": {"giftcards": {"on": True, "qty": 1, "rate": "500"}},
        "proposalDisc": "10",
    }
    out = os.path.join(_HERE, "test_proposal.pptx")
    with open(out, "wb") as f:
        f.write(build_pptx(sample))
    print(f"Written → {out}")
